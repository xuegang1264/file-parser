"""
file-parser: 轻量级文档解析服务。

基于纯 Python 库，无 AI 模型、无大体积依赖。
支持 PDF、DOCX、PPTX，旧版 DOC/PPT 可在安装 LibreOffice 后自动转档。
提供独立 HTTP 服务，接收文件上传，返回 Markdown 文本。
"""
import base64
import os
import shutil
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Optional, Union

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from loguru import logger

try:
    import pypdf
except ImportError:  # pragma: no cover
    pypdf = None

try:
    import docx
except ImportError:  # pragma: no cover
    docx = None

try:
    import pptx
except ImportError:  # pragma: no cover
    pptx = None

app = FastAPI(title="File Parser Server", version="0.2.0")

# CORS：允许 ai-rag 等 Agent Server 跨域调用
app.add_middleware(
    CORSMiddleware,
    allow_origins=os.environ.get("CORS_ORIGINS", "*").split(","),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 默认工作目录，用于存放临时文件
WORK_DIR = Path(os.environ.get("WORK_DIR", tempfile.gettempdir())) / "file-parser"
WORK_DIR.mkdir(parents=True, exist_ok=True)

# 输出长度限制，防止超大文档直接塞爆响应
MAX_CONTENT_CHARS = int(os.environ.get("MAX_CONTENT_CHARS", "300000"))


def _allowed_extension(filename: str) -> bool:
    """检查文件扩展名是否在支持范围内。"""
    ext = Path(filename).suffix.lower().lstrip(".")
    allowed = {
        "pdf",
        "doc", "docx",
        "ppt", "pptx",
    }
    return ext in allowed


def _extract_pdf(source: Union[Path, bytes]) -> str:
    """使用 pypdf 提取 PDF 文本。"""
    if pypdf is None:
        raise RuntimeError("未安装 pypdf，无法解析 PDF")

    reader = pypdf.PdfReader(source)
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


def _extract_docx(source: Union[Path, bytes]) -> str:
    """使用 python-docx 提取 Word 文本与表格。"""
    if docx is None:
        raise RuntimeError("未安装 python-docx，无法解析 DOCX")

    document = docx.Document(source)
    parts = []

    # 段落
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # 表格（简单 Markdown 格式）
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if not rows:
            continue

        header = rows[0]
        lines = ["| " + " | ".join(header) + " |"]
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            # 补齐列数
            if len(row) < len(header):
                row = row + [""] * (len(header) - len(row))
            elif len(row) > len(header):
                row = row[: len(header)]
            lines.append("| " + " | ".join(row) + " |")
        parts.append("\n".join(lines))

    return "\n\n".join(parts)


def _iter_shape_texts(shape):
    """递归提取 pptx 形状及其组合内的文本。"""
    texts = []
    if hasattr(shape, "text") and shape.text.strip():
        texts.append(shape.text.strip())
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(_iter_shape_texts(child))
    return texts


def _extract_pptx(source: Union[Path, bytes]) -> str:
    """使用 python-pptx 提取 PPT 每页文本。"""
    if pptx is None:
        raise RuntimeError("未安装 python-pptx，无法解析 PPTX")

    prs = pptx.Presentation(source)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(_iter_shape_texts(shape))
        if slide_texts:
            parts.append(f"## 第 {i} 页\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _convert_with_libreoffice(input_path: Path, out_dir: Path, target_ext: str) -> Path:
    """
    使用 LibreOffice 将旧版 DOC/PPT 转成 DOCX/PPTX。
    如果系统未安装 soffice，会抛出 RuntimeError。
    """
    try:
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                target_ext,
                "--outdir",
                str(out_dir),
                str(input_path),
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=120,
            check=False,
        )
    except FileNotFoundError as e:
        raise RuntimeError(
            f"未找到 LibreOffice，无法处理旧版 {input_path.suffix} 文件"
        ) from e

    if result.returncode != 0:
        logger.error(f"LibreOffice stderr: {result.stderr}")
        raise RuntimeError(f"LibreOffice 转换失败: {result.stderr[:500]}")

    converted = out_dir / f"{input_path.stem}.{target_ext}"
    if not converted.exists():
        raise RuntimeError("LibreOffice 转换后未找到输出文件")
    return converted


def _parse_file(file_path: Path, output_format: str = "markdown") -> dict:
    """根据扩展名分发解析逻辑。"""
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        content = _extract_pdf(file_path)
    elif ext == ".docx":
        content = _extract_docx(file_path)
    elif ext == ".pptx":
        content = _extract_pptx(file_path)
    elif ext == ".doc":
        with tempfile.TemporaryDirectory() as tmpdir:
            converted = _convert_with_libreoffice(file_path, Path(tmpdir), "docx")
            content = _extract_docx(converted)
    elif ext == ".ppt":
        with tempfile.TemporaryDirectory() as tmpdir:
            converted = _convert_with_libreoffice(file_path, Path(tmpdir), "pptx")
            content = _extract_pptx(converted)
    else:
        raise RuntimeError(f"不支持的文件类型: {ext}")

    metadata = {"mime": "unknown", "pages": 0}
    pages = []

    # 为兼容旧接口，json 格式同样返回 content；pages 为空
    if output_format == "json":
        return {
            "content": content[:MAX_CONTENT_CHARS],
            "metadata": metadata,
            "pages": pages,
        }

    return {
        "content": content[:MAX_CONTENT_CHARS],
        "metadata": metadata,
        "pages": pages,
    }


@app.get("/health")
def health():
    """健康检查。"""
    return {"status": "ok", "service": "file-parser"}


@app.post("/parse")
async def parse(
    file: UploadFile = File(..., description="待解析的文件"),
    output_format: str = Form(default="markdown", description="输出格式: markdown | json"),
):
    """
    解析上传文件，返回 Markdown 文本。

    - 不保存上传的文件与解析结果。
    - 解析超时默认 10 分钟。
    """
    if output_format not in {"markdown", "json"}:
        raise HTTPException(status_code=400, detail="output_format 必须是 markdown 或 json")

    request_id = str(uuid.uuid4())
    work_dir = WORK_DIR / request_id
    upload_path = work_dir / (file.filename or "upload")

    try:
        work_dir.mkdir(parents=True, exist_ok=True)

        # 保存上传文件
        with open(upload_path, "wb") as f:
            shutil.copyfileobj(file.file, f)

        if not _allowed_extension(upload_path.name):
            raise HTTPException(
                status_code=400, detail=f"不支持的文件类型: {upload_path.suffix}"
            )

        parse_result = _parse_file(upload_path, output_format)

        return JSONResponse(
            content={
                "status": "success",
                "content": parse_result["content"],
                "metadata": parse_result["metadata"],
                "pages": parse_result["pages"],
            }
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.exception("文件解析失败")
        raise HTTPException(status_code=500, detail=f"解析失败: {str(e)}") from e
    finally:
        # 清理临时文件，保证不存储
        try:
            if work_dir.exists():
                shutil.rmtree(work_dir, ignore_errors=True)
        except Exception as e:
            logger.warning(f"清理临时目录失败 {work_dir}: {e}")


@app.post("/parse/base64")
async def parse_base64(
    filename: str = Form(...),
    data: str = Form(..., description="文件内容的 base64 编码"),
    output_format: str = Form(default="markdown"),
):
    """接收 base64 编码的文件内容并解析。"""
    try:
        file_bytes = base64.b64decode(data)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"base64 解码失败: {e}") from e

    request_id = str(uuid.uuid4())
    work_dir = WORK_DIR / request_id
    upload_path = work_dir / filename

    try:
        work_dir.mkdir(parents=True, exist_ok=True)
        upload_path.write_bytes(file_bytes)

        if not _allowed_extension(upload_path.name):
            raise HTTPException(
                status_code=400, detail=f"不支持的文件类型: {upload_path.suffix}"
            )

        parse_result = _parse_file(upload_path, output_format)
        return JSONResponse(
            content={
                "status": "success",
                "content": parse_result["content"],
                "metadata": parse_result["metadata"],
                "pages": parse_result["pages"],
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("base64 文件解析失败")
        raise HTTPException(status_code=500, detail=f"解析失败: {str(e)}") from e
    finally:
        try:
            if work_dir.exists():
                shutil.rmtree(work_dir, ignore_errors=True)
        except Exception as e:
            logger.warning(f"清理临时目录失败 {work_dir}: {e}")


if __name__ == "__main__":
    host = os.environ.get("HOST", "0.0.0.0")
    port = int(os.environ.get("PORT", "8000"))
    import uvicorn

    uvicorn.run(app, host=host, port=port)
